import os
import fitz
import tabula
import camelot
import pandas as pd
import io
from flask import Flask, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from docx import Document
from pdf2docx import Converter
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from io import StringIO
from bs4 import BeautifulSoup
from flask_cors import CORS

app = Flask(__name__)
CORS = (app)

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
    return response

# Set absolute paths for uploads and converted folders
base_dir = os.path.dirname(os.path.abspath(__file__))
upload_folder = os.path.join(base_dir, 'uploads')
converted_folder = os.path.join(base_dir, 'converted')

# Create subdirectories for each conversion type
text_dir = os.path.join(converted_folder, 'text')
word_dir = os.path.join(converted_folder, 'word')
img_dir = os.path.join(converted_folder, 'images')
excel_dir = os.path.join(converted_folder, 'excel')
ppt_dir = os.path.join(converted_folder, 'ppt')

# Ensure the directories exist
os.makedirs(upload_folder, exist_ok=True)
os.makedirs(text_dir, exist_ok=True)
os.makedirs(word_dir, exist_ok=True)
os.makedirs(img_dir, exist_ok=True)
os.makedirs(excel_dir, exist_ok=True)
os.makedirs(ppt_dir, exist_ok=True)

app.config['UPLOAD_FOLDER'] = upload_folder
app.config['CONVERTED_FOLDER'] = converted_folder
app.config['ALLOWED_EXTENSIONS'] = {'pdf'}

# Helper function to check file extensions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# PDF to Text conversion
from pytesseract import image_to_string

def pdf_to_text(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page_number, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if not page_text.strip():
            print(f"Falling back to OCR for page {page_number}")
            # Use fitz to convert page to an image
            doc = fitz.open(pdf_path)
            pix = doc[page_number].get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes()))
            page_text = image_to_string(img)
        text += page_text
    text_file = os.path.join(text_dir, f'{os.path.splitext(os.path.basename(pdf_path))[0]}.txt')
    with open(text_file, 'w', encoding='utf-8') as f:
        f.write(text)
    return text_file

# PDF to Word conversion with formatting preservation
def pdf_to_word(pdf_path):
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"File {pdf_path} does not exist.")
    
    # Convert PDF to Word using pdf2docx
    word_file = os.path.join(word_dir, f'{os.path.splitext(os.path.basename(pdf_path))[0]}.docx')
    cv = Converter(pdf_path)
    cv.convert(word_file, start=0, end=None)  # Convert the entire PDF
    cv.close()  # Close the converter after conversion

    return word_file


# PDF to Image conversion

def pdf_to_image(pdf_path):
    doc = fitz.open(pdf_path)
    img_paths = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=300)
        img_path = os.path.join(img_dir, f'{os.path.splitext(os.path.basename(pdf_path))[0]}_{page_num+1}.png')
        pix.save(img_path)
        img_paths.append(img_path)
    return img_paths

def pdf_to_excel(pdf_path):
    """
    Convert a PDF file to Excel format using Camelot and Tabula as a fallback.
    """
    try:
        import camelot  # Ensure Camelot is installed
    except ImportError:
        raise ImportError("Camelot library is not installed. Install it using 'pip install camelot-py[cv]'.")

    excel_files = []
    try:
        # Primary attempt: Camelot
        tables = camelot.read_pdf(pdf_path, pages='all', strip_text='\n', split_text=True)

        if not tables or len(tables) == 0:
            raise ValueError("No tables found in the PDF using Camelot.")

        for idx, table in enumerate(tables):
            # Convert table to a pandas DataFrame
            df = table.df

            # Save each table to an Excel file
            excel_filename = f"{os.path.splitext(os.path.basename(pdf_path))[0]}_camelot_table_{idx + 1}.xlsx"
            excel_file_path = os.path.join(excel_dir, excel_filename)
            df.to_excel(excel_file_path, index=False, header=True)
            excel_files.append(excel_filename)

        return excel_files

    except Exception as e:
        print(f"Camelot failed: {str(e)}")
        print("Attempting fallback with Tabula...")

        # Fallback: Tabula
        try:
            import tabula
        except ImportError:
            raise ImportError("Tabula library is not installed. Install it using 'pip install tabula-py'.")

        try:
            tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)

            if not tables or len(tables) == 0:
                raise ValueError("No tables found using Tabula either.")

            for idx, table in enumerate(tables):
                # Save Tabula DataFrame to an Excel file
                excel_filename = f"{os.path.splitext(os.path.basename(pdf_path))[0]}_tabula_table_{idx + 1}.xlsx"
                excel_file_path = os.path.join(excel_dir, excel_filename)
                table.to_excel(excel_file_path, index=False, header=True)
                excel_files.append(excel_filename)

        except Exception as fallback_error:
            raise ValueError(f"Both Camelot and Tabula failed: {str(fallback_error)}")

    return excel_files

# PDF to PPT conversion
def pdf_to_ppt(pdf_path):
    ppt_file = os.path.join(ppt_dir, f"{os.path.splitext(os.path.basename(pdf_path))[0]}.pptx")
    presentation = Presentation()
    reader = PdfReader(pdf_path)

    for page_num, page in enumerate(reader.pages):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout

        # Extract text from the PDF page
        page_text = page.extract_text() or "No extractable content found"

        # Use OCR if text extraction fails
        if not page_text.strip():
            print(f"Falling back to OCR for page {page_num}")
            doc = fitz.open(pdf_path)
            pix = doc[page_num].get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes()))
            page_text = image_to_string(img)

        # Add text to the slide
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        frame = textbox.text_frame
        frame.word_wrap = True
        for paragraph_text in page_text.split('\n'):
            paragraph = frame.add_paragraph()
            paragraph.text = paragraph_text
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Arial"

        # Add page image as a background or additional content
        doc = fitz.open(pdf_path)
        pix = doc[page_num].get_pixmap(dpi=150)
        img_path = os.path.join(img_dir, f"{os.path.splitext(os.path.basename(pdf_path))[0]}_{page_num + 1}.png")
        pix.save(img_path)
        
        left = Inches(0.5)
        top = Inches(6)
        slide.shapes.add_picture(img_path, left, top, width=Inches(9))

    presentation.save(ppt_file)
    return ppt_file

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'message': 'No file part'}), 400
    file = request.files['file']
    if file and allowed_file(file.filename):
        # Sanitize filename by replacing spaces with dashes
        filename = secure_filename(file.filename).replace(' ', '-')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        print(f"Saving file to: {file_path}")
        file.save(file_path)
        return jsonify({'message': 'File uploaded successfully', 'filename': filename}), 200
    return jsonify({'message': 'Invalid file format'}), 400

@app.route('/converted', methods=['POST'])
def convert_file():
    data = request.json
    filename = secure_filename(data.get('filename')).replace(' ', '-')  # Sanitize the filename
    conversion_type = data.get('conversion_type')

    if not filename or not conversion_type:
        return jsonify({'message': 'Missing filename or conversion type'}), 400

    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)

    if not os.path.exists(file_path):
        return jsonify({'message': 'File not found'}), 404

    try:
        converted_files = []
        print(f"Processing file: {file_path}")

        if conversion_type == 'text':
            text_file = pdf_to_text(file_path)
            converted_files.append({'type': 'text', 'filename': os.path.basename(text_file)})

        elif conversion_type == 'word':
            word_file = pdf_to_word(file_path)
            converted_files.append({'type': 'word', 'filename': os.path.basename(word_file)})

        elif conversion_type == 'image':
            img_paths = pdf_to_image(file_path)
            converted_files.extend([{'type': 'image', 'filename': os.path.basename(img)} for img in img_paths])

        elif conversion_type == 'excel':
            excel_files = pdf_to_excel(file_path)
            converted_files.extend([{'type': 'excel', 'filename': excel_file} for excel_file in excel_files])
        
        elif conversion_type == 'ppt':
            ppt_file = pdf_to_ppt(file_path)
            converted_files.append({'type': 'ppt', 'filename': os.path.basename(ppt_file)})

        else:
            return jsonify({'message': 'Invalid conversion type'}), 400

        return jsonify({'converted_files': converted_files}), 200

    except Exception as e:
        print(f"Error: {str(e)}")
        return jsonify({'message': f"Error converting file: {str(e)}"}), 500

@app.route('/download/<conversion_type>/<filename>', methods=['GET'])
def download_file(conversion_type, filename):
    folder_map = {
        'text': text_dir,
        'word': word_dir,
        'image': img_dir,
        'excel': excel_dir,
        'ppt': ppt_dir,
    }

    if conversion_type not in folder_map:
        return jsonify({'message': 'Invalid conversion type'}), 400

    folder = folder_map[conversion_type]

    # Sanitize filename by replacing spaces with dashes
    sanitized_filename = secure_filename(filename).replace(' ', '-')

    file_path = os.path.join(folder, sanitized_filename)
    if not os.path.exists(file_path):
        return jsonify({'message': 'File not found'}), 404

    return send_from_directory(folder, sanitized_filename)


if __name__ == '__main__':
    app.run(debug=True)
